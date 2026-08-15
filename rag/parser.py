"""
Stage 1: DOCUMENT -> TEXT / TABLES / IMAGES

Splits a PDF or PPTX into three raw streams, each tagged with the page (or
slide) number it came from so citations can point back to it later.
This stage never calls an external API - extraction of raw structure is
always local, it's just parsing a file format.
"""
import os
import io
import uuid
from dataclasses import dataclass, field
from typing import List

import fitz  # PyMuPDF
import pdfplumber
from PIL import Image
from pptx import Presentation


@dataclass
class RawText:
    doc_name: str
    page: int
    text: str


@dataclass
class RawTable:
    doc_name: str
    page: int
    rows: List[List[str]]
    table_index: int


@dataclass
class RawImage:
    doc_name: str
    page: int
    image_id: str
    path: str
    image_index: int


@dataclass
class ParsedDocument:
    doc_name: str
    texts: List[RawText] = field(default_factory=list)
    tables: List[RawTable] = field(default_factory=list)
    images: List[RawImage] = field(default_factory=list)


def _save_image(pil_img: Image.Image, images_dir: str) -> str:
    os.makedirs(images_dir, exist_ok=True)
    image_id = f"{uuid.uuid4().hex[:12]}.png"
    path = os.path.join(images_dir, image_id)
    if pil_img.mode in ("RGBA", "P"):
        pil_img = pil_img.convert("RGB")
    pil_img.save(path, format="PNG")
    return path


def parse_pdf(file_path: str, doc_name: str, images_dir: str) -> ParsedDocument:
    doc = ParsedDocument(doc_name=doc_name)

    # Text + tables via pdfplumber (best for layout-aware text and grid tables)
    with pdfplumber.open(file_path) as pdf:
        for page_idx, page in enumerate(pdf.pages, start=1):
            text = (page.extract_text() or "").strip()
            if text:
                doc.texts.append(RawText(doc_name=doc_name, page=page_idx, text=text))

            tables = page.extract_tables() or []
            for t_idx, table in enumerate(tables):
                cleaned = [[(cell or "").strip() for cell in row] for row in table]
                cleaned = [row for row in cleaned if any(cell for cell in row)]
                if cleaned:
                    doc.tables.append(
                        RawTable(doc_name=doc_name, page=page_idx, rows=cleaned, table_index=t_idx)
                    )

    # Images via PyMuPDF (handles embedded raster images reliably)
    fitz_doc = fitz.open(file_path)
    for page_idx in range(len(fitz_doc)):
        page = fitz_doc[page_idx]
        for img_idx, img in enumerate(page.get_images(full=True)):
            xref = img[0]
            try:
                base_image = fitz_doc.extract_image(xref)
                image_bytes = base_image["image"]
                pil_img = Image.open(io.BytesIO(image_bytes))
                if min(pil_img.size) < 40:
                    continue  # skip tiny decorative icons/bullets
                path = _save_image(pil_img, images_dir)
                doc.images.append(
                    RawImage(
                        doc_name=doc_name,
                        page=page_idx + 1,
                        image_id=os.path.basename(path),
                        path=path,
                        image_index=img_idx,
                    )
                )
            except Exception:
                continue
    fitz_doc.close()
    return doc


def parse_pptx(file_path: str, doc_name: str, images_dir: str) -> ParsedDocument:
    doc = ParsedDocument(doc_name=doc_name)
    prs = Presentation(file_path)

    for slide_idx, slide in enumerate(prs.slides, start=1):
        slide_text_parts = []
        table_counter = 0
        image_counter = 0

        for shape in slide.shapes:
            # Text frames (titles, bullets, text boxes)
            if shape.has_text_frame and shape.text_frame.text.strip():
                slide_text_parts.append(shape.text_frame.text.strip())

            # Native PPTX tables
            if shape.has_table:
                tbl = shape.table
                rows = []
                for row in tbl.rows:
                    rows.append([cell.text.strip() for cell in row.cells])
                if rows:
                    doc.tables.append(
                        RawTable(doc_name=doc_name, page=slide_idx, rows=rows, table_index=table_counter)
                    )
                    table_counter += 1

            # Pictures
            if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                try:
                    image_bytes = shape.image.blob
                    pil_img = Image.open(io.BytesIO(image_bytes))
                    if min(pil_img.size) < 40:
                        continue
                    path = _save_image(pil_img, images_dir)
                    doc.images.append(
                        RawImage(
                            doc_name=doc_name,
                            page=slide_idx,
                            image_id=os.path.basename(path),
                            path=path,
                            image_index=image_counter,
                        )
                    )
                    image_counter += 1
                except Exception:
                    continue

        # Speaker notes are useful context too
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                slide_text_parts.append(f"[Speaker notes] {notes}")

        if slide_text_parts:
            doc.texts.append(RawText(doc_name=doc_name, page=slide_idx, text="\n".join(slide_text_parts)))

    return doc


def parse_document(file_path: str, doc_name: str, images_dir: str) -> ParsedDocument:
    ext = os.path.splitext(file_path)[1].lower()
    if ext == ".pdf":
        return parse_pdf(file_path, doc_name, images_dir)
    elif ext in (".pptx",):
        return parse_pptx(file_path, doc_name, images_dir)
    else:
        raise ValueError(f"Unsupported file type: {ext}. Use .pdf or .pptx")
