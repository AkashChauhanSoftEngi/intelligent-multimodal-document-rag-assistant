"""Generates sample_data/Q4_Revenue_Deck.pptx - a small multimodal demo deck
with text, a native table, and a chart image."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

OUT = os.path.join(os.path.dirname(__file__), "..", "sample_data", "Q4_Revenue_Deck.pptx")
CHART = os.path.join(os.path.dirname(__file__), "..", "sample_data", "revenue_chart.png")

prs = Presentation()
blank = prs.slide_layouts[6]
title_layout = prs.slide_layouts[0]

# Slide 1 - Title
slide = prs.slides.add_slide(title_layout)
slide.shapes.title.text = "Q4 FY2025 Revenue Review"
slide.placeholders[1].text = "Prepared for the Board of Directors — January 2026"

# Slide 2 - Text summary
slide = prs.slides.add_slide(blank)
box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(9), Inches(1))
box.text_frame.text = "Q4 Highlights"
box.text_frame.paragraphs[0].font.size = Pt(28)
box.text_frame.paragraphs[0].font.bold = True

body = slide.shapes.add_textbox(Inches(0.6), Inches(1.4), Inches(9), Inches(4))
tf = body.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Full-year revenue reached $51.0M, up 18% from $43.2M in FY2024."
p.font.size = Pt(16)
for line in [
    "Cloud Services grew 31.7% year-over-year, the fastest of any segment.",
    "Gross margin expanded to 62%, up from 58% in FY2024.",
    "Asia-Pacific revenue grew 31% year-over-year on new enterprise contracts.",
]:
    para = tf.add_paragraph()
    para.text = line
    para.font.size = Pt(16)

speaker_notes = slide.notes_slide.notes_text_frame
speaker_notes.text = "Emphasize that Cloud Services is now the primary growth driver going into FY2026."

# Slide 3 - Table
slide = prs.slides.add_slide(blank)
box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(9), Inches(0.6))
box.text_frame.text = "Revenue by Segment"
box.text_frame.paragraphs[0].font.size = Pt(24)
box.text_frame.paragraphs[0].font.bold = True

rows, cols = 5, 4
table_shape = slide.shapes.add_table(rows, cols, Inches(0.6), Inches(1.2), Inches(9), Inches(2.6))
table = table_shape.table
headers = ["Segment", "FY2024 ($M)", "FY2025 ($M)", "YoY Growth"]
data = [
    ["Enterprise Software", "18.5", "23.1", "+24.9%"],
    ["Cloud Services", "14.2", "18.7", "+31.7%"],
    ["Professional Services", "7.1", "6.4", "-9.9%"],
    ["Support & Maintenance", "3.4", "2.8", "-17.6%"],
]
for c, h in enumerate(headers):
    cell = table.cell(0, c)
    cell.text = h
    cell.text_frame.paragraphs[0].font.bold = True
    cell.text_frame.paragraphs[0].font.size = Pt(13)
for r, row in enumerate(data, start=1):
    for c, val in enumerate(row):
        cell = table.cell(r, c)
        cell.text = val
        cell.text_frame.paragraphs[0].font.size = Pt(12)

# Slide 4 - Chart image
slide = prs.slides.add_slide(blank)
box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(9), Inches(0.6))
box.text_frame.text = "Figure 8 — Revenue by Year"
box.text_frame.paragraphs[0].font.size = Pt(24)
box.text_frame.paragraphs[0].font.bold = True
slide.shapes.add_picture(CHART, Inches(1.7), Inches(1.3), width=Inches(6.5))

prs.save(OUT)
print(f"Wrote {OUT}")
